import os
import tempfile
import pandas as pd
import streamlit as st
import zipfile
import tempfile
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.documents import Document as LangchainDocument

from pptx import Presentation
from docx import Document as DocxDocument

# -------------------------
# App Config
# -------------------------
st.set_page_config(page_title="ModdleBot", layout="wide")
st.title("📘MoodleBot– RAG App")

# -------------------------
# Constants
# -------------------------
CHROMA_DIR = "chroma_db"

# -------------------------
# Sidebar Controls
# -------------------------
st.sidebar.header("⚙️ Settings")

max_tokens = st.sidebar.slider("Max Tokens", 50, 300, 100)
temperature = st.sidebar.slider("Temperature", 0.0, 1.0, 0.7)
topic = st.sidebar.text_input("Current Topic", "General")

# -------------------------
# Load Embeddings & Vector Store
# -------------------------
@st.cache_resource
def load_vector_store():
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )
    return Chroma(
        embedding_function=embeddings,
        persist_directory=CHROMA_DIR
    )

vector_store = load_vector_store()

# -------------------------
# Extract Text from PPTX
# -------------------------
def extract_ppt_text(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# -------------------------
# Extract Text from DOCX
# -------------------------
def extract_docx_text(path):
    doc = DocxDocument(path)
    return "\n".join([para.text for para in doc.paragraphs])

# -------------------------
# Process ZIP File
# -------------------------
def process_zip_file(zip_file):

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )

    chunks = []

    with tempfile.TemporaryDirectory() as tmpdir:

        # Save uploaded ZIP temporarily
        zip_path = os.path.join(tmpdir, "uploaded.zip")
        with open(zip_path, "wb") as f:
            f.write(zip_file.read())

        # Extract ZIP
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(tmpdir)

        # Walk through extracted files
        for root, _, files in os.walk(tmpdir):
            for file in files:
                file_path = os.path.join(root, file)

                try:
                    # TXT
                    if file.endswith(".txt"):
                        loader = TextLoader(file_path, encoding="utf-8")
                        docs = loader.load()

                    # PDF
                    elif file.endswith(".pdf"):
                        loader = PyPDFLoader(file_path)
                        docs = loader.load()

                    # PPTX
                    elif file.endswith(".pptx"):
                        text = extract_ppt_text(file_path)
                        docs = [LangchainDocument(
                            page_content=text,
                            metadata={"source": file}
                        )]

                    # DOCX
                    elif file.endswith(".docx"):
                        text = extract_docx_text(file_path)
                        docs = [LangchainDocument(
                            page_content=text,
                            metadata={"source": file}
                        )]

                    else:
                        continue

                    # Split into chunks
                    split_docs = splitter.split_documents(docs)
                    chunks.extend(split_docs)

                except Exception as e:
                    st.warning(f"Error processing {file}: {e}")

    # Store embeddings
    if chunks:
        vector_store.add_documents(chunks)
        #vector_store.persist()
        st.sidebar.success(f"✅ {len(chunks)} chunks added to vector store")
    else:
        st.sidebar.warning("No valid documents found in ZIP")

# -------------------------
# Sidebar: ZIP Upload
# -------------------------
st.sidebar.subheader("📂 Upload Learning Materials (ZIP)")

uploaded_zip = st.sidebar.file_uploader(
    "Upload ZIP file (txt, pdf, pptx, docx inside)",
    type=["zip"]
)

if st.sidebar.button("📥 Process ZIP"):
    if uploaded_zip:
        process_zip_file(uploaded_zip)
    else:
        st.sidebar.warning("Upload a ZIP file first")

# -------------------------
# Student Classification (NORMAL)
# -------------------------
def classify_student(df):
    numeric = df.select_dtypes(include=["int64", "float64"])
    avg_score = numeric.mean(axis=1).iloc[0]

    if avg_score < 2.5:
        return "Below Average"
    elif avg_score < 4:
        return "Average"
    else:
        return "Above Average"

# -------------------------
# Sidebar: Excel Upload
# -------------------------
st.sidebar.subheader("📊 Student Performance")

uploaded_excel = st.sidebar.file_uploader(
    "Upload Marks Excel (.xlsx)",
    type=["xlsx"]
)

learning_level = "Average"

if uploaded_excel:
    df = pd.read_excel(uploaded_excel)
    learning_level = classify_student(df)

    st.sidebar.success(f"Student Level: {learning_level}")

# -------------------------
# Load Gemini Model
# -------------------------
@st.cache_resource
def load_model():
    return ChatGoogleGenerativeAI(
        model="models/gemini-2.5-flash",
        api_key="AIzaSyDi13U7Zi1j1i-0KOR5wNl1XRaDv6mI6b8",
        temperature=temperature,
        #max_output_tokens=max_tokens
    )

qa_llm = load_model()

# -------------------------
# Chat History
# -------------------------
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# -------------------------
# User Question
# -------------------------
question = st.text_input("❓ Ask a question")

if st.button("Ask") and question:
    retriever = vector_store.as_retriever(
    search_type="mmr",   # better diversity
    search_kwargs={
        "k": 8,          # retrieve more chunks
        "fetch_k": 20    # search from larger pool
    })
    docs = retriever.invoke(question)
    print(docs)
    context = "\n".join(d.page_content for d in docs)
    print("ihg")
    print(context)

    prompt = f"""
You are an AI tutor.

STRICT RULES:
- Answer ONLY using the information present in the provided context.
- DO NOT use outside knowledge.
- If the answer is not found in the context, say:
  "The provided documents do not contain enough information to answer this question."

Student Learning Level: {learning_level}

Guidelines based on learning level:
- Below Average: Use very simple words, short sentences, step-by-step explanation, and a small example if available.
- Average: Use clear explanation with key points and brief examples.
- Above Average: Use concise, technical explanation with proper terminology.

Context:
{context}

Question:
{question}

Answer:
"""

    response = qa_llm.invoke(prompt)
    answer = response.content.strip()

    st.session_state.chat_history.append({
        "user": question,
        "bot": answer
    })

# -------------------------
# Display Chat
# -------------------------
st.subheader("💬 Chat History")

for chat in reversed(st.session_state.chat_history):
    st.markdown(f"**You:** {chat['user']}")
    st.markdown(f"**AI:** {chat['bot']}")
    st.divider()
